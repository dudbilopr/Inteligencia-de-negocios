import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Colores UNAB
UNAB_BLUE = RGBColor(0, 56, 101) # #003865
UNAB_ORANGE = RGBColor(232, 92, 11) # #E85C0B
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)

def apply_theme(slide, bg_color=UNAB_BLUE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

def format_text(shape, text, size_pt, bold=False, rgb=WHITE, align=PP_ALIGN.LEFT):
    shape.text = text
    for p in shape.text_frame.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.color.rgb = rgb
            run.font.name = 'Century Gothic'

def add_bullet(tf, text, level=0, bold=False, rgb=WHITE, size=24):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb
    p.font.name = 'Century Gothic'
    # Adjust line spacing to prevent overflow
    p.space_after = Pt(10)

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]

    # --- SLIDE 1: Title (High Impact) ---
    slide = prs.slides.add_slide(blank_layout)
    apply_theme(slide, UNAB_BLUE)
    tb1 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
    format_text(tb1, "Inteligencia de Negocios", size_pt=60, bold=True, rgb=UNAB_ORANGE, align=PP_ALIGN.CENTER)
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1))
    format_text(tb2, "Clase 1: Introducción, Evolución y Conceptos Básicos\nDocente: Dudbil Olvasada Pabon Riaño", size_pt=30, bold=False, rgb=WHITE, align=PP_ALIGN.CENTER)

    def add_content_slide(title_text):
        sl = prs.slides.add_slide(blank_layout)
        apply_theme(sl, LIGHT_GRAY)
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(1))
        format_text(tb, title_text, size_pt=40, bold=True, rgb=UNAB_BLUE)
        shape = sl.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = UNAB_ORANGE
        shape.line.color.rgb = UNAB_ORANGE
        body = sl.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.8))
        tf = body.text_frame
        tf.word_wrap = True
        return sl, tf

    # --- SLIDE 2: Agenda ---
    sl, tf = add_content_slide("Agenda de la Sesión")
    add_bullet(tf, "1. Acuerdos y Reglas de Clase", bold=True, rgb=UNAB_ORANGE, size=28)
    add_bullet(tf, "2. Estrategia EMI (Warm-up en Inglés)", bold=True, rgb=UNAB_ORANGE, size=28)
    add_bullet(tf, "3. ¿Qué es Inteligencia de Negocios (BI)?", bold=True, rgb=UNAB_ORANGE, size=28)
    add_bullet(tf, "4. Glosario Técnico y Fundamentos Analíticos", bold=True, rgb=UNAB_ORANGE, size=28)
    add_bullet(tf, "5. Storytelling with Data (Knaflic)", bold=True, rgb=UNAB_ORANGE, size=28)
    add_bullet(tf, "6. Taller Práctico: Pair Programming", bold=True, rgb=UNAB_ORANGE, size=28)
    add_bullet(tf, "7. Conclusiones y Bibliografía", bold=True, rgb=UNAB_ORANGE, size=28)

    # --- SLIDE 3: Acuerdos ---
    sl, tf = add_content_slide("Acuerdos de Clase")
    add_bullet(tf, "Horario Oficial: Viernes de 10:00 AM a 1:00 PM.", bold=True, rgb=UNAB_BLUE, size=26)
    add_bullet(tf, "Puntualidad: Llegar a tiempo garantiza no perder conceptos clave.", level=1, rgb=UNAB_BLUE, size=22)
    add_bullet(tf, "Descansos Cognitivos: Haremos una pausa a las 11:30 AM para evitar fatiga técnica.", level=1, rgb=UNAB_BLUE, size=22)
    add_bullet(tf, "Estrategia EMI: 10 a 15 minutos de inmersión en inglés técnico por sesión.", level=1, rgb=UNAB_BLUE, size=22)

    # --- SLIDE 4: EMI Warm up ---
    sl, tf = add_content_slide("🇬🇧 EMI Strategy: Warm-Up")
    add_bullet(tf, "Let's activate our English and BI Vocabulary!", bold=True, rgb=UNAB_BLUE, size=32)
    add_bullet(tf, "\nActividad Interactiva:", level=0, size=28, rgb=UNAB_BLUE)
    add_bullet(tf, "Abran el archivo quiz_semana_1.html en sus navegadores.", level=1, bold=True, rgb=UNAB_ORANGE, size=26)
    add_bullet(tf, "Vamos a emparejar los conceptos (Dashboard, KPI, Data-Driven) con su definición.", level=1, rgb=UNAB_BLUE, size=22)

    # --- SLIDE 5: Que es BI ---
    sl, tf = add_content_slide("¿Qué es la Inteligencia de Negocios?")
    add_bullet(tf, "Definición Institucional (Sharda et al., 2020):", bold=True, rgb=UNAB_BLUE, size=28)
    add_bullet(tf, "Conjunto de estrategias, procesos y tecnologías enfocadas en la creación de conocimiento mediante análisis de datos.", level=1, rgb=UNAB_BLUE, size=22)
    add_bullet(tf, "\nEl rol de la Analítica:", bold=True, rgb=UNAB_ORANGE, size=28)
    add_bullet(tf, "- Descriptiva: ¿Qué sucedió? (Reportes históricos)", level=1, rgb=UNAB_BLUE, size=22)
    add_bullet(tf, "- Predictiva: ¿Qué pasará? (Modelos predictivos, Machine Learning)", level=1, rgb=UNAB_BLUE, size=22)
    add_bullet(tf, "- Prescriptiva: ¿Qué debemos hacer? (Optimización, Simulación)", level=1, rgb=UNAB_BLUE, size=22)

    # --- SLIDE 6: Glosario Técnico ---
    sl, tf = add_content_slide("Glosario Técnico de BI")
    add_bullet(tf, "ETL (Extract, Transform, Load):", bold=True, rgb=UNAB_BLUE, size=26)
    add_bullet(tf, "Extraer datos, transformarlos al formato analítico y cargarlos. El 80% del trabajo en BI ocurre aquí.", level=1, rgb=UNAB_BLUE, size=20)
    add_bullet(tf, "OLTP vs OLAP:", bold=True, rgb=UNAB_BLUE, size=26)
    add_bullet(tf, "OLTP: Transacciones rápidas diarias (Ej. Caja Registradora). Altamente normalizado.", level=1, rgb=UNAB_BLUE, size=20)
    add_bullet(tf, "OLAP: Consultas analíticas complejas y agregaciones (Cubo de Datos).", level=1, rgb=UNAB_BLUE, size=20)
    add_bullet(tf, "Data Warehouse & Data Mart:", bold=True, rgb=UNAB_BLUE, size=26)
    add_bullet(tf, "Data Warehouse: Repositorio centralizado de datos de toda la corporación.", level=1, rgb=UNAB_BLUE, size=20)
    add_bullet(tf, "Data Mart: Una subdivisión del Warehouse para un solo departamento (Ej. Finanzas).", level=1, rgb=UNAB_BLUE, size=20)

    # --- SLIDE 7: Storytelling ---
    sl, tf = add_content_slide("Storytelling con Datos")
    add_bullet(tf, "Principios de Diseño Visual (Knaflic, 2015):", bold=True, rgb=UNAB_ORANGE, size=28)
    add_bullet(tf, "1. Eliminar el desorden (Declutter):", bold=True, rgb=UNAB_BLUE, size=24)
    add_bullet(tf, "Elimina bordes de gráficos, gridlines de fondo, colores innecesarios y sombras.", level=1, rgb=UNAB_BLUE, size=20)
    add_bullet(tf, "2. Dirigir la Atención (Principios Gestalt):", bold=True, rgb=UNAB_BLUE, size=24)
    add_bullet(tf, "Usa colores vibrantes (naranja) SOLO para los datos clave. El resto en tonos neutros (gris/azul).", level=1, rgb=UNAB_BLUE, size=20)
    add_bullet(tf, "3. Evitar el 3D a toda costa:", bold=True, rgb=UNAB_BLUE, size=24)
    add_bullet(tf, "Los gráficos 3D (especialmente Pie Charts) distorsionan la percepción geométrica del cerebro.", level=1, rgb=UNAB_BLUE, size=20)

    # --- SLIDE 8: Imagen Simulada (Dashboard Clean) ---
    sl = prs.slides.add_slide(blank_layout)
    apply_theme(sl, LIGHT_GRAY)
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(1))
    format_text(tb, "Simulación: Dashboard de Alto Impacto (Knaflic)", size_pt=35, bold=True, rgb=UNAB_BLUE)
    shape = sl.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = UNAB_ORANGE
    shape.line.color.rgb = UNAB_ORANGE
    
    img_path = 'dashboard_mockup.png'
    if os.path.exists(img_path):
        sl.shapes.add_picture(img_path, Inches(2), Inches(1.5), width=Inches(9.3))
    else:
        tb_err = sl.shapes.add_textbox(Inches(2), Inches(3), Inches(8), Inches(1))
        format_text(tb_err, "Imagen simulada del Dashboard no encontrada.", size_pt=24, bold=True, rgb=UNAB_BLUE)

    # --- SLIDE 9: Taller ---
    sl, tf = add_content_slide("Taller: Pair Programming")
    add_bullet(tf, "Análisis Crítico de Dashboards en el Mundo Real", bold=True, rgb=UNAB_ORANGE, size=30)
    add_bullet(tf, "1. En parejas, busquen en internet 2 Dashboards (Uno saturado/malo y uno limpio/bueno).", level=1, rgb=UNAB_BLUE, size=24)
    add_bullet(tf, "2. Analicen qué principios de Gestalt y 'Decluttering' se violan en el malo.", level=1, rgb=UNAB_BLUE, size=24)
    add_bullet(tf, "3. Analicen la jerarquía visual del Dashboard bueno.", level=1, rgb=UNAB_BLUE, size=24)
    add_bullet(tf, "4. Redacten un resumen en formato Markdown y súbanlo a su repositorio de GitHub.", level=1, rgb=UNAB_BLUE, size=24)
    add_bullet(tf, "Nota: Usen LLMNotebook/ChatGPT como asistente crítico para la revisión.", level=1, bold=True, rgb=UNAB_BLUE, size=24)

    # --- SLIDE 10: Conclusiones ---
    sl, tf = add_content_slide("Conclusiones de la Sesión")
    add_bullet(tf, "El BI no es solo tecnología, es una herramienta estratégica:", bold=True, rgb=UNAB_BLUE, size=28)
    add_bullet(tf, "El verdadero valor del BI está en pasar de los datos descriptivos (pasado) a las decisiones prescriptivas (acciones).", level=1, rgb=UNAB_BLUE, size=24)
    add_bullet(tf, "El diseño visual importa tanto como el modelo de datos:", bold=True, rgb=UNAB_BLUE, size=28)
    add_bullet(tf, "Un modelo perfecto (Data Warehouse robusto) fracasa si el Dashboard es ilegible (Exceso de Clutter).", level=1, rgb=UNAB_BLUE, size=24)

    # --- SLIDE 11: Bibliografía ---
    sl, tf = add_content_slide("Bibliografía y Referencias")
    add_bullet(tf, "Lecturas Obligatorias:", bold=True, rgb=UNAB_ORANGE, size=28)
    add_bullet(tf, "1. Sharda, R., Delen, D., & Turban, E. (2020). Business Intelligence, Analytics, and Data Science: A Managerial Perspective (4th Edition). Pearson.", level=1, rgb=UNAB_BLUE, size=22)
    add_bullet(tf, "2. Knaflic, C. N. (2015). Storytelling with Data: A Data Visualization Guide for Business Professionals. Wiley.", level=1, rgb=UNAB_BLUE, size=22)
    add_bullet(tf, "\nRecursos Adicionales:", bold=True, rgb=UNAB_ORANGE, size=28)
    add_bullet(tf, "Repositorio del curso: https://github.com/dudbil/Inteligencia-de-negocios-UNAB", level=1, rgb=UNAB_BLUE, size=22)

    prs.save('diapositivas_semana_1.pptx')
    print("Presentación PPTX generada con éxito con imágenes, conclusiones y bibliografía.")

if __name__ == '__main__':
    main()
